#!/usr/bin/env python3
"""Build the matrix-free CVFEM talk as a polished .pptx.

Narrative: title (decode the jargon) -> motivation (why the operator is the wall)
-> assembly wall -> matrix-free + sum-factorization -> accuracy -> scaling/trillion
-> outlook. Figures from make_matfree_figs.py. Current numbers (620M/GPU -> ~1,700 GPU,
~0.38 TDOF/s). Starter on the blank master -- apply your institutional template after.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
BLUE   = RGBColor(0x1F, 0x4E, 0x79)
LBLUE  = RGBColor(0x2E, 0x75, 0xB6)
ORANGE = RGBColor(0xC0, 0x5A, 0x10)
GREY   = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0x8A, 0x8A, 0x8A)
DARK   = RGBColor(0x20, 0x20, 0x20)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BAND   = RGBColor(0x10, 0x2A, 0x43)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
FOOT = "UniDistance.ch   ·   USI Euler Inst.   ·   HSLU Luzern   ·   CSCS   ·   PASC"

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp

def textbox(s, x, y, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor: tf.vertical_anchor = anchor
    return tf

def run(p, text, size, color=DARK, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    return r

def header(s, text):
    tf = textbox(s, 0.55, 0.30, 12.2, 0.85)
    run(tf.paragraphs[0], text, 28, BLUE, bold=True)
    rect(s, 0.6, 1.12, 4.2, 0.045, ORANGE)        # accent rule
    footer(s)

def footer(s):
    tf = textbox(s, 0.55, 7.05, 12.2, 0.35)
    run(tf.paragraphs[0], FOOT, 9, LGREY)

def bullets(s, items, x=0.55, y=1.45, w=6.3, h=5.3, size=16, gap=8):
    tf = textbox(s, x, y, w, h)
    for i, it in enumerate(items):
        txt, lvl = it[0], it[1]
        color = it[2] if len(it) > 2 else None
        bold  = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4 if lvl == 0 else gap)
        if lvl == 0:        # sub-head
            run(p, txt, size + 3, BLUE, bold=True)
        elif lvl == 1:      # bullet
            run(p, "•  ", size, ORANGE, bold=True); run(p, txt, size, color or DARK, bold=bold)
        else:               # sub-bullet
            p.level = 1
            run(p, "–  ", size - 1, LGREY); run(p, txt, size - 1, color or GREY)

def image(s, name, x=7.15, y=1.6, w=5.7):
    path = os.path.join(HERE, name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

# ============================ 1. TITLE ============================
s = slide()
rect(s, 0, 0, 13.333, 2.55, BAND)
tf = textbox(s, 0.7, 0.55, 12.0, 1.9, anchor=MSO_ANCHOR.MIDDLE)
run(tf.paragraphs[0], "Matrix-Free High-Order CVFEM at Scale", 38, WHITE, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(6)
run(p, "Solving fluid dynamics on billions of points — without ever storing the matrix", 19, RGBColor(0xCF,0xE0,0xF0), italic=True)
# decoder lines
tf = textbox(s, 0.9, 3.05, 11.6, 2.7)
dec = [
    ("Matrix-free", "never assemble or store A — recompute y = A·x on the fly, the inner loop of every solve"),
    ("High-order", "~570× more accuracy per unknown, at the same throughput as low order"),
    ("At scale", "40 billion unknowns today on 64 GPUs — a trillion-class operator, Gordon-Bell density"),
]
for i, (k, v) in enumerate(dec):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    run(p, "▸ ", 18, ORANGE, bold=True)
    run(p, k + ":  ", 18, BLUE, bold=True)
    run(p, v, 17, DARK)
tf = textbox(s, 0.9, 6.35, 11.6, 0.9)
run(tf.paragraphs[0], "Daniel Ganellari et al.        [affiliations]        PASC / Gordon Bell  ·  2026", 14, GREY)

# ============================ 2. MOTIVATION ============================
s = slide()
header(s, "Why: the operator apply is the wall")
bullets(s, [
    ("High-fidelity CFD needs two things at once", 0),
    ("Accuracy — resolve the physics (turbulence, boundary layers, pumps, turbines)", 1),
    ("Scale — billions of unknowns to capture a real geometry", 1),
    ("Every implicit / iterative solve spends almost all its time in one operation", 0),
    ("y = A·x  — the operator apply, run hundreds of times per solve", 1, ORANGE, True),
    ("High order is how you buy accuracy per unknown — but the price is memory", 0),
    ("The assembled high-order matrix explodes and starves on bandwidth", 1),
    ("This work", 0),
    ("A high-order CVFEM operator that is cheap to store, fast to apply,", 1),
    ("bit-exact distributed, and scales toward a trillion DOF on GPUs", 1),
], w=12.2, h=5.4)

# ============================ 3. THE WALL ============================
s = slide()
header(s, "The assembly wall: high order is too big to store")
bullets(s, [
    ("Storage by mode — and where each one wins", 0),
    ("Assembled 7-nnz CSR: ~84 B/DOF, flat in p (CVFEM stencil)", 1),
    ("p=1: assembled (84) is LEANER than matrix-free (PA 278 / MF 185 B/DOF)", 1, ORANGE, True),
    ("Matrix-free wins storage only at high order — crossover p ≈ 2", 1),
    ("The plot", 0),
    ("MF-recompute (green) drops fast with p; PA (blue) stays moderate; assembled (orange) flat", 1),
])
image(s, "fig_memory.png", y=1.7)

# ============================ 4. MATRIX-FREE + SUMFAC ============================
s = slide()
header(s, "Matrix-free + sum-factorization: high order, for free")
bullets(s, [
    ("Store only the geometry, recompute the apply", 0),
    ("One per-element metric (d_G); no global matrix, no coordinates in the kernel", 1),
    ("Trilinos-free: cstone octree + own CSR + cuSPARSE + Hypre, GPU-native", 1),
    ("Sum-factorization: high order is cheap to store", 0),
    ("O(p⁴) work; recompute storage drops ~240× from p=1→8 — high order is the lean regime", 1, ORANGE, True),
    ("The apply was the wall — shared-memory-instruction bound, not bandwidth", 0),
    ("ncu: LSU/shared pipe ~99% busy, HBM idle; FP64 tensor cores don't help here", 1, GREY),
    ("Register + warp-shuffle apply (Nek-style): contractions in registers, not shared", 1),
    ("~2× at the peak — ~6–13 GDOF/s/GPU, peak p=3; A·1 bit-exact at every order", 1, ORANGE, True),
], w=6.6)
image(s, "fig_throughput.png", x=7.4, y=1.7, w=5.5)

# ============================ 5. ACCURACY ============================
s = slide()
header(s, "Accuracy: high order is worth it — and it's correct")
bullets(s, [
    ("Same unknowns, far less error", 0),
    ("~570× lower error per DOF, p=1 → p=4 (at ~2197 DOF: 2.0e-2 vs 3.5e-5)", 1),
    ("~p+1 convergence order", 1),
    ("Distributed, bit-for-bit", 0),
    ("A·1 = machine-zero: 1e-18 on 4 ranks (patch test)", 1, ORANGE, True),
    ("Device halo bit-identical to host — distribution reproduces single-rank exactly", 1),
])
image(s, "fig_convergence.png", y=1.7)

# ============================ 6. SCALING / TRILLION ============================
s = slide()
header(s, "Scaling: 40 billion DOF today, a trillion-class operator")
bullets(s, [
    ("Weak-scales without noticing it's distributed", 0),
    ("40B DOF on 64 GH200 (p=4): apply ~100% weak-scaling, full matvec ~90%", 1),
    ("~0.38 TDOF/s sustained; one full matvec ~0.1 s regardless of scale", 1, ORANGE, True),
    ("Communication self-resolves: 22% → 4% as each GPU's block grows", 1),
    ("Toward a trillion", 0),
    ("620M DOF/GPU ceiling (p=4) → 10¹² projected on ~1,700 GH200 (~430 nodes)", 1),
    ("Trillion-CLASS operator — not a trillion result yet (projection, not a run)", 1, GREY),
], w=6.7, size=15)
image(s, "fig_trillion.png", x=7.5, y=1.9, w=5.4)

# ============================ 7. OUTLOOK ============================
s = slide()
header(s, "Outlook: from operator to solve, and to adaptivity")
bullets(s, [
    ("The solve (next)", 0),
    ("FlexGMRES + low-order-refined AMG, on top of the validated operator", 1),
    ("Genuinely unstructured at scale", 0),
    ("Conforming 10⁹-DOF meshes via uniform refinement of an unstructured base", 1),
    ("Differentiator", 0),
    ("Native unstructured AMR on the cstone octree — not structured-AMR overset", 1, ORANGE, True),
    ("Takeaway", 0),
    ("A trustworthy, bit-exact, trillion-class high-order CVFEM operator on GPUs —", 1, BLUE, True),
    ("the foundation the solve and the adaptivity build on", 1, BLUE, True),
], w=12.2, h=5.4)

out = os.path.join(os.path.dirname(HERE), "pasc_matrixfree.pptx")
prs.save(out)
print("wrote", out, "(%d slides)" % len(prs.slides._sldIdLst))
